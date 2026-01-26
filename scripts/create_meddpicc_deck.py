#!/usr/bin/env python3
"""
Create a modern MEDDPICC Sales Training PowerPoint deck
with The CRO Report branding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# CRO Report Brand Colors
NAVY = RGBColor(0x1E, 0x3A, 0x5F)  # #1e3a5f
NAVY_LIGHT = RGBColor(0x2D, 0x4A, 0x6F)  # #2d4a6f
NAVY_DARK = RGBColor(0x15, 0x2D, 0x4A)  # #152d4a
GOLD = RGBColor(0xD9, 0x77, 0x06)  # #d97706
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100 = RGBColor(0xF1, 0xF5, 0xF9)  # #f1f5f9
GRAY_500 = RGBColor(0x64, 0x74, 0x8B)  # #64748b
GRAY_600 = RGBColor(0x47, 0x55, 0x69)  # #475569

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def set_background_color(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=NAVY, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide with navy background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_background_color(slide, NAVY)

    # Gold accent bar at top
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), GOLD)

    # Title
    add_text_box(slide, Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.5),
                 title, font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.8),
                     subtitle, font_size=24, color=GRAY_100, align=PP_ALIGN.CENTER)

    # Footer with branding
    add_text_box(slide, Inches(0.75), Inches(6.8), Inches(5), Inches(0.4),
                 "THE CRO REPORT", font_size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    return slide

def create_section_slide(prs, section_title, section_number=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, NAVY_DARK)

    # Large section number
    if section_number:
        add_text_box(slide, Inches(0.75), Inches(1.5), Inches(2), Inches(1.5),
                     section_number, font_size=120, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Section title
    add_text_box(slide, Inches(0.75), Inches(3.5), Inches(11), Inches(1.5),
                 section_title, font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Gold accent line
    add_rectangle(slide, Inches(0.75), Inches(5.2), Inches(3), Inches(0.08), GOLD)

    return slide

def create_content_slide(prs, title, bullets, has_icon=False, icon_text=""):
    """Create a standard content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, WHITE)

    # Navy header bar
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), NAVY)

    # Title on header
    add_text_box(slide, Inches(0.75), Inches(0.35), Inches(11), Inches(0.8),
                 title, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Gold accent under header
    add_rectangle(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.06), GOLD)

    # Content area
    content_top = Inches(1.8)
    bullet_text = "\n".join([f"• {b}" for b in bullets])
    add_text_box(slide, Inches(0.75), content_top, Inches(11.5), Inches(5),
                 bullet_text, font_size=22, color=GRAY_600, align=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(0.75), Inches(7), Inches(5), Inches(0.3),
                 "THE CRO REPORT", font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide

def create_meddpicc_letter_slide(prs, letter, title, content_bullets, is_current=True):
    """Create a MEDDPICC letter focus slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, WHITE)

    # Navy header bar
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), NAVY)

    # Gold accent
    add_rectangle(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.06), GOLD)

    # Large letter indicator
    letter_box = add_rounded_rectangle(slide, Inches(0.5), Inches(0.25), Inches(0.9), Inches(0.9), GOLD)
    # Add letter text on top
    add_text_box(slide, Inches(0.5), Inches(0.28), Inches(0.9), Inches(0.9),
                 letter, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.6), Inches(0.35), Inches(10), Inches(0.8),
                 title, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # MEDDPICC indicator bar on right
    letters = ['M', 'E', 'D', 'D', 'P', 'I', 'C', 'C']
    letter_idx = {'M': 0, 'E': 1, 'D1': 2, 'D2': 3, 'P': 4, 'I': 5, 'C1': 6, 'C2': 7}

    for i, l in enumerate(letters):
        box_top = Inches(1.6 + i * 0.65)
        is_active = (l == letter) or (letter.startswith('D') and l == 'D' and i in [2, 3]) or \
                    (letter.startswith('C') and l == 'C' and i in [6, 7])

        color = GOLD if is_active else GRAY_100
        text_color = WHITE if is_active else GRAY_500

        box = add_rounded_rectangle(slide, Inches(12.2), box_top, Inches(0.6), Inches(0.55), color)
        add_text_box(slide, Inches(12.2), box_top + Inches(0.05), Inches(0.6), Inches(0.5),
                     l, font_size=20, bold=True, color=text_color, align=PP_ALIGN.CENTER)

    # Content bullets
    bullet_text = "\n\n".join([f"• {b}" for b in content_bullets])
    add_text_box(slide, Inches(0.75), Inches(1.8), Inches(11), Inches(5),
                 bullet_text, font_size=20, color=GRAY_600, align=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(0.75), Inches(7), Inches(5), Inches(0.3),
                 "THE CRO REPORT", font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide

def create_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Create a two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, WHITE)

    # Navy header bar
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), NAVY)
    add_text_box(slide, Inches(0.75), Inches(0.35), Inches(11), Inches(0.8),
                 title, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rectangle(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.06), GOLD)

    # Left column
    left_box = add_rounded_rectangle(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5), GRAY_100)
    add_text_box(slide, Inches(0.75), Inches(1.9), Inches(5.3), Inches(0.6),
                 left_title, font_size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    left_content = "\n".join([f"• {b}" for b in left_bullets])
    add_text_box(slide, Inches(0.75), Inches(2.6), Inches(5.3), Inches(3.8),
                 left_content, font_size=18, color=GRAY_600, align=PP_ALIGN.LEFT)

    # Right column
    right_box = add_rounded_rectangle(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(5), GRAY_100)
    add_text_box(slide, Inches(7.05), Inches(1.9), Inches(5.3), Inches(0.6),
                 right_title, font_size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    right_content = "\n".join([f"• {b}" for b in right_bullets])
    add_text_box(slide, Inches(7.05), Inches(2.6), Inches(5.3), Inches(3.8),
                 right_content, font_size=18, color=GRAY_600, align=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(0.75), Inches(7), Inches(5), Inches(0.3),
                 "THE CRO REPORT", font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide

def create_quote_slide(prs, quote, attribution=""):
    """Create an impactful quote slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, NAVY_DARK)

    # Large quote marks
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(2), Inches(1.5),
                 '"', font_size=200, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Quote text
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(3),
                 quote, font_size=32, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Attribution
    if attribution:
        add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                     f"— {attribution}", font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

    return slide

def create_stage_slide(prs, stage_num, stage_name, activities, outcomes):
    """Create a sales stage slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, WHITE)

    # Navy header bar
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), NAVY)

    # Stage number badge
    add_rounded_rectangle(slide, Inches(0.5), Inches(0.25), Inches(1.2), Inches(0.9), GOLD)
    add_text_box(slide, Inches(0.5), Inches(0.32), Inches(1.2), Inches(0.8),
                 f"STAGE {stage_num}", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Stage name
    add_text_box(slide, Inches(1.9), Inches(0.35), Inches(10), Inches(0.8),
                 stage_name, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_rectangle(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.06), GOLD)

    # Activities column
    add_text_box(slide, Inches(0.75), Inches(1.6), Inches(5.5), Inches(0.5),
                 "ACTIVITIES", font_size=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    activities_text = "\n".join([f"• {a}" for a in activities])
    add_text_box(slide, Inches(0.75), Inches(2.1), Inches(5.5), Inches(4.5),
                 activities_text, font_size=16, color=GRAY_600, align=PP_ALIGN.LEFT)

    # Outcomes column
    add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5),
                 "CUSTOMER OUTCOMES", font_size=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    outcomes_text = "\n".join([f"✓ {o}" for o in outcomes])
    add_text_box(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4.5),
                 outcomes_text, font_size=16, color=GRAY_600, align=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(0.75), Inches(7), Inches(5), Inches(0.3),
                 "THE CRO REPORT", font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    return slide

def create_meddpicc_overview_slide(prs):
    """Create the MEDDPICC overview grid slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, NAVY_DARK)

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "The MEDDPICC Framework", font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rectangle(slide, Inches(0.5), Inches(1.0), Inches(2), Inches(0.06), GOLD)

    # MEDDPICC items
    items = [
        ("M", "Metrics", "What is the economic impact of our deal?"),
        ("E", "Economic Buyer", "Who has P&L responsibility and budget authority?"),
        ("D", "Decision Criteria", "What is their technical, vendor, and financial criteria?"),
        ("D", "Decision Process", "Which steps are required in their buying process?"),
        ("P", "Identify Pain", "What's motivating them to do anything?"),
        ("I", "Implicate Pain", "What's the impact if they don't act?"),
        ("C", "Champion", "Who has the power and conviction to sell on your behalf?"),
        ("C", "Competition", "Who are we competing against?"),
    ]

    for i, (letter, name, desc) in enumerate(items):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 6.2)
        top = Inches(1.4 + row * 1.4)

        # Letter badge
        add_rounded_rectangle(slide, left, top, Inches(0.7), Inches(0.7), GOLD)
        add_text_box(slide, left, top + Inches(0.05), Inches(0.7), Inches(0.65),
                     letter, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Name and description
        add_text_box(slide, left + Inches(0.9), top, Inches(5), Inches(0.5),
                     name, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_text_box(slide, left + Inches(0.9), top + Inches(0.5), Inches(5), Inches(0.6),
                     desc, font_size=14, color=GRAY_100, align=PP_ALIGN.LEFT)

    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # === SLIDE 1: Title ===
    create_title_slide(prs,
                       "MEDDPICC Sales Mastery",
                       "A Strategic Framework for Sales Excellence")

    # === SLIDE 2: Table of Contents ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, WHITE)
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), NAVY)
    add_text_box(slide, Inches(0.75), Inches(0.35), Inches(11), Inches(0.8),
                 "What We'll Cover", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rectangle(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.06), GOLD)

    toc_items = [
        ("01", "Understanding MEDDPICC", "The qualifying framework that drives forecast accuracy"),
        ("02", "Deep Dive: Each Letter", "Metrics, Economic Buyer, Decision Criteria & more"),
        ("03", "The Three Whys", "Why do anything? Why us? Why now?"),
        ("04", "Champions", "Identifying, developing, and leveraging your champion"),
        ("05", "Sales Stages", "Mapping MEDDPICC to your sales process"),
    ]

    for i, (num, title, desc) in enumerate(toc_items):
        top = Inches(1.7 + i * 1.05)
        add_text_box(slide, Inches(0.75), top, Inches(0.8), Inches(0.5),
                     num, font_size=32, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
        add_text_box(slide, Inches(1.6), top, Inches(4), Inches(0.5),
                     title, font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        add_text_box(slide, Inches(1.6), top + Inches(0.45), Inches(10), Inches(0.5),
                     desc, font_size=16, color=GRAY_500, align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(0.75), Inches(7), Inches(5), Inches(0.3),
                 "THE CRO REPORT", font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # === SECTION 1: Understanding MEDDPICC ===
    create_section_slide(prs, "Understanding MEDDPICC", "01")

    # What is MEDDPICC
    create_content_slide(prs, "What is MEDDPICC?", [
        "A qualification framework to assess deal health and knowledge gaps",
        "Helps identify blind spots in your opportunities",
        "Creates a common language across your sales organization",
        "Maintains forecast accuracy through systematic qualification",
        "Not a checklist — it's a journey throughout the entire sales process"
    ])

    # MEDDPICC Overview Grid
    create_meddpicc_overview_slide(prs)

    # Benefits
    create_content_slide(prs, "Why MEDDPICC Matters", [
        "Assess deal health objectively — no more gut feelings",
        "Identify blind spots before they derail your deal",
        "Speak the same language with your team and leadership",
        "Maintain forecast accuracy with data-driven qualification",
        "Prepare strategically for every prospect interaction"
    ])

    # Quote slide
    create_quote_slide(prs,
                       "MEDDPICC is a journey, not an event or checklist. Use it throughout the entire sales process.",
                       "")

    # === SECTION 2: The Three Whys ===
    create_section_slide(prs, "The Three Whys", "02")

    create_content_slide(prs, "Every Deal Must Answer Three Questions", [
        "WHY DO ANYTHING? — What's driving them to change?",
        "WHY US? — What makes us the right choice?",
        "WHY NOW? — What creates urgency to act?",
        "",
        "If you can't answer all three, your deal is at risk."
    ])

    # === SECTION 3: Deep Dive ===
    create_section_slide(prs, "MEDDPICC Deep Dive", "03")

    # Pain
    create_meddpicc_letter_slide(prs, "P", "Identify Pain", [
        "Current Situation: What is their current approach (or lack of one)?",
        "Uncovered Pain: What business and personal frustrations exist?",
        "Negative Consequences: What happens if they don't address this?",
        "",
        "No delta, no deal. Pain drives action."
    ])

    # Metrics
    create_meddpicc_letter_slide(prs, "M", "Metrics", [
        "How will the prospect measure business impact from your solution?",
        "What economic outcomes matter most to them?",
        "Quantify the value — make it tangible and measurable",
        "",
        "Metrics create the business case for change."
    ])

    # Champion
    create_meddpicc_letter_slide(prs, "C", "Champion", [
        "Power & Influence: Someone who grants access to the Economic Buyer",
        "Personal Win: They have a vested interest in your success",
        "Sells for You: They articulate Why Do Anything, Why Now, Why Us",
        "",
        "A champion is not the same as a coach. Know the difference."
    ])

    # Champion vs Coach
    create_two_column_slide(prs, "Champion vs. Coach",
        "CHAMPION", [
            "Has power and influence",
            "Has a personal stake in your win",
            "Actively sells when you're not there",
            "Can articulate the 3 Whys",
            "Returns your calls proactively"
        ],
        "COACH", [
            "Provides information",
            "May not have decision power",
            "Helps you navigate, not sell",
            "May be friendly but passive",
            "Limited influence with EB"
        ])

    # Decision Criteria & Competition
    create_meddpicc_letter_slide(prs, "D", "Decision Criteria", [
        "Technical Criteria: What capabilities do they require?",
        "Vendor Criteria: What do they need from a partner?",
        "Financial Criteria: What's the budget and ROI expectation?",
        "",
        "Bake your differentiators into their criteria."
    ])

    # Competition strategies
    create_content_slide(prs, "Competitive Strategies", [
        "DIRECT — When you have overwhelming advantage (3:1). Default strategy.",
        "DIVIDE — Fragment the competition. 'Beach Head' or 'Trojan Horse' approach.",
        "INDIRECT — When competition is too strong. Change the ground by leveraging new buying influences.",
        "DELAY — When losing. Slow the deal to change MEDDPICC dynamics.",
        "",
        "Know your enemy's position. Attack where unprepared."
    ])

    # Economic Buyer
    create_meddpicc_letter_slide(prs, "E", "Economic Buyer", [
        "Has discretionary use of funds and veto power",
        "Owns the bottom-line impact of the decision",
        "Can say YES when everyone else says NO",
        "",
        "Validate: Who is the final authorizer in the paper process?"
    ])

    # Decision Process
    create_meddpicc_letter_slide(prs, "D", "Decision Process", [
        "Selection Process: How will they decide what to do?",
        "Contract Process: What steps to complete the vendor process?",
        "",
        "Map every step. Own the timeline. No surprises."
    ])

    # === SECTION 4: Champions ===
    create_section_slide(prs, "Mastering Champions", "04")

    create_content_slide(prs, "The Champion Definition", [
        "POWER & INFLUENCE — They command the room",
        "PERSONAL WIN — They have something to gain",
        "SELLS FOR YOU — They advocate when you're not there",
        "",
        "You need evidence of all three."
    ])

    create_content_slide(prs, "How to Spot a Champion", [
        "They answer the tough questions with authority",
        "They're not afraid to say 'I don't know'",
        "They've led transformational projects before",
        "They're hard to reach (high demand)",
        "New to the organization with something to prove"
    ])

    create_content_slide(prs, "Developing Your Champion", [
        "Add value to them — be honest, have integrity",
        "Don't be a 'salesperson' — be a trusted advisor",
        "Predict the future — spot problems before they do",
        "Provide information they can reuse as their own",
        "Walk them through MEDDPICC, help them prepare their Champion Deck"
    ])

    create_content_slide(prs, "Testing Your Champion", [
        "Can they articulate the 3 Whys without your guidance?",
        "Can they give you access to the EB and run meetings for you?",
        "Do they provide competitive intelligence proactively?",
        "Do they return your calls? Do they call YOU?",
        "Play the negative: 'I'm not sure we can get this done...' — see how they respond"
    ])

    create_quote_slide(prs, "No Champion, No Deal.\nDevelop Multiple Champions.", "")

    # === SECTION 5: Sales Stages ===
    create_section_slide(prs, "Sales Stage Execution", "05")

    create_content_slide(prs, "Sales Methodology Principles", [
        "What is the pain? (Technical, Business, Personal)",
        "Who must be involved? Who has power?",
        "Where is the risk?",
        "",
        "We are measured by executed contracts, not advancing opportunities."
    ])

    # Stage slides
    create_stage_slide(prs, "0", "Identify an Opportunity",
        ["BDR schedules intro call", "Create opportunity once scheduled",
         "Identify: Pain, Economic Buyer, Timing", "Build alignment on case for change"],
        ["Shares pain and why they want to learn more", "Agrees to discovery meeting"])

    create_stage_slide(prs, "1", "Determine Problem & Impact",
        ["Conduct discovery call", "Identify metrics customer wants to improve",
         "Map decision process and paper process", "Identify potential champions",
         "Confirm pricing, timeline, competitors"],
        ["Attends discovery meeting", "Confirms evaluation & buying decision path"])

    create_stage_slide(prs, "2", "Validate Benefits & Value",
        ["Present demo to key stakeholders", "Introduce Mutual Action Plan (MAP)",
         "Influence business & technical criteria", "Demonstrate value vs. competitors",
         "Test and validate champion has influence"],
        ["Attends demo", "Agrees to rollout discussion with EB"])

    create_stage_slide(prs, "3", "Confirm Value with Power",
        ["Discuss rollout and implementation", "Meet EB or person who influences EB",
         "Present business case to stakeholders", "Align on commercial structure",
         "Deliver references"],
        ["Attends rollout meeting", "Confirms vision lock & rollout strategy",
         "Starts InfoSec & Legal review"])

    create_stage_slide(prs, "4", "Negotiate & Execute",
        ["Deliver final order form to EB/Champion", "Confirm vendor of choice",
         "Verify remaining concerns are met", "Approve InfoSec & Legal",
         "Execute contract"],
        ["InfoSec & Legal approved", "Contract is signed"])

    create_stage_slide(prs, "5", "Closed Won",
        ["Introduce Customer Success", "Joint AE & CS handoff call",
         "Ensure EB is satisfied with success metrics"],
        ["EB & stakeholders attend kickoff call", "Partnership begins"])

    # === CLOSING ===
    create_title_slide(prs, "Execute with MEDDPICC", "Qualify Better. Forecast Accurately. Close More Deals.")

    # Save
    output_path = "/Users/rome/Documents/croreport/MEDDPICC_Sales_CRO_Report.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    main()
